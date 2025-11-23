#!/usr/bin/env python3
"""
Generate PowerPoint presentation for MSC & Co EIC Accelerator Application
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

# Color palette
COLORS = {
    'primary_blue': RGBColor(*hex_to_rgb('#667eea')),
    'purple': RGBColor(*hex_to_rgb('#764ba2')),
    'pink': RGBColor(*hex_to_rgb('#f857a6')),
    'dark': RGBColor(*hex_to_rgb('#0a0e27')),
    'green': RGBColor(*hex_to_rgb('#51cf66')),
    'red': RGBColor(*hex_to_rgb('#ff6b6b')),
    'gold': RGBColor(*hex_to_rgb('#ffd700')),
    'white': RGBColor(255, 255, 255),
    'light_gray': RGBColor(200, 200, 200),
}

def create_presentation():
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # SLIDE 1: COVER
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_dark_gradient_background(slide1)
    add_slide_number(slide1, "(01)")

    # Logo
    logo = slide1.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(8), Inches(1))
    logo_frame = logo.text_frame
    logo_frame.text = "MSC & Co"
    p = logo_frame.paragraphs[0]
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']

    # Subtitle
    subtitle = slide1.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(10), Inches(0.3))
    subtitle_frame = subtitle.text_frame
    subtitle_frame.text = "AUDIOMSC LTD • Company No. 13250829"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['light_gray']

    # Main title
    title = slide1.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(14), Inches(2))
    title_frame = title.text_frame
    title_frame.text = "AI-Native Climate-Positive\nMusic Distribution Ecosystem"
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.font.size = Pt(68)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']
    p.line_spacing = 1.1

    # Innovation cards (simplified as text boxes)
    innovations = [
        ("🌍", "Real-time carbon attribution in music streaming"),
        ("⛓️", "Blockchain-verified royalty transparency"),
        ("🤖", "Conversational AI for distribution automation")
    ]

    x_start = 0.7
    card_width = 4.5
    y = 4.5
    for i, (icon, text) in enumerate(innovations):
        x = x_start + (i * (card_width + 0.3))
        add_innovation_card(slide1, x, y, card_width, icon, text)

    # Premium badge
    badge = slide1.shapes.add_textbox(Inches(3.5), Inches(7.2), Inches(9), Inches(0.8))
    badge_frame = badge.text_frame
    badge_frame.text = "€1.8M EIC Accelerator Grant Request"
    p = badge_frame.paragraphs[0]
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # SLIDE 2: MARKET OPPORTUNITY
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_gradient_background(slide2)
    add_slide_number(slide2, "(02)")
    add_title(slide2, "Market Opportunity: €28.6B Industry Crisis", "Independent Artists Face Systemic Exploitation")

    # Stats in 2x2 grid
    stats = [
        ("💰", "€8.5B", "Independent Artist Market", "40% of streaming revenue"),
        ("🎵", "8.5M", "Artists Trapped", "Between Exploitation & Impossibility"),
        ("📈", "15%", "YoY Growth", "Independent Sector (2x Overall Market)"),
        ("🌡️", "300M", "Tonnes CO₂/Year", "ZERO Platforms Track Emissions")
    ]

    grid_x = [1, 8.5]
    grid_y = [2.5, 5]
    idx = 0
    for row in range(2):
        for col in range(2):
            icon, value, label, sublabel = stats[idx]
            add_stat_card(slide2, grid_x[col], grid_y[row], icon, value, label, sublabel)
            idx += 1

    # Callout box
    callout = slide1.shapes.add_textbox(Inches(1), Inches(7.8), Inches(14), Inches(0.8))
    callout_frame = callout.text_frame
    callout_frame.text = "Legacy platforms: 15-30% commission FOREVER • No carbon tracking • No blockchain transparency • Zero AI automation"
    p = callout_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # SLIDE 3: TECHNOLOGY MOAT
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_gradient_background(slide3)
    add_slide_number(slide3, "(03)")
    add_title(slide3, "18-24 Month Technology Lead", "First-Mover Advantage in Three Breakthrough Areas")

    features = [
        ("🚀", "First MCP-Native Platform", "Built on Model Context Protocol (Anthropic, 2024). Competitors stuck on 20-year-old architecture."),
        ("🤖", "181 MCP Tools Live", "Complete distribution automation. Competitors have ZERO conversational AI capabilities."),
        ("✅", "85% Complete", "Production-ready platform. Launch in 90 days post-funding, not 18-24 months."),
        ("🎯", "Apollo Intelligence", "Live conversational AI. Hit prediction, fraud detection, automated KYC compliance.")
    ]

    idx = 0
    for row in range(2):
        for col in range(2):
            icon, title, text = features[idx]
            add_feature_card(slide3, grid_x[col], grid_y[row], icon, title, text)
            idx += 1

    # Comparison panel
    comp_y = 7.5
    add_comparison_text(slide3, 1.5, comp_y, "Traditional Platforms", "4 HOURS", "15 forms • Manual • No guidance", COLORS['red'])
    add_text_box(slide3, 7, comp_y, 2, 0.5, "VS", Pt(48), COLORS['white'], True, PP_ALIGN.CENTER)
    add_comparison_text(slide3, 10, comp_y, "MSC & Co (MCP-Native)", "5 MINUTES", "One conversation • Apollo AI • Automated", COLORS['green'])

    # SLIDE 4: SUSTAINABILITY
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_gradient_background(slide4)
    add_slide_number(slide4, "(04)")
    add_title(slide4, "First Climate-Positive Music Platform", "Leading the Music Industry's Green Transformation")

    sustainability_features = [
        ("🌍", "Real-Time Carbon Tracking", "Per-stream emissions using DIMPACT 2024 methodology. No competitor has this capability."),
        ("🤝", "Earth/Percent Partnership", "Founded by Massive Attack & Brian Eno. Automatic royalty donations to climate action projects."),
        ("♻️", "Carbon Offset Marketplace", "5 verified providers integrated: Gold Standard, Verra, Ecologi, Pachama, Climeworks.")
    ]

    for i, (icon, title, text) in enumerate(sustainability_features):
        x = 1 + (i * 4.8)
        add_feature_card(slide4, x, 2.5, icon, title, text, width=4.5)

    # Impact banner
    add_impact_banner(slide4, 1.5, 6.5, 13, "5-YEAR COMPETITIVE MOAT",
                     "Competitors would need complete platform rebuild to match sustainability features", COLORS['green'])

    # SLIDE 5: BLOCKCHAIN
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_gradient_background(slide5)
    add_slide_number(slide5, "(05)")
    add_title(slide5, "Blockchain-Verified Royalty Transparency", "Cryptographic Proof > Trust")

    blockchain_features = [
        ("🔒", "IMMUTABLE\nRoyalty Records on Polygon", "Every payment recorded on blockchain. Artists independently verify earnings without trusting platform."),
        ("💎", "99.97%\nCheaper Than Ethereum", "$0.005 per transaction using Polygon. Ethereum would cost $200+ per verification."),
        ("🤝", "SMART\nContract Revenue Splits", "Automated distribution to collaborators. No manual calculations. Transparent waterfall visualization."),
        ("✅", "PROOF\nNot Trust Required", "Artists get cryptographic proof of earnings. Audit trail viewable by anyone. Zero competitors have this.")
    ]

    idx = 0
    for row in range(2):
        for col in range(2):
            icon, title, text = blockchain_features[idx]
            add_feature_card(slide5, grid_x[col], grid_y[row], icon, title, text)
            idx += 1

    # Quote banner
    add_quote_banner(slide5, 1.5, 7.8, 13, "The music industry has a $2.5B trust problem. We solve it with cryptographic PROOF.")

    # SLIDE 6: EUROPEAN VALUE
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_gradient_background(slide6)
    add_slide_number(slide6, "(06)")
    add_title(slide6, "European Champion: Digital Sovereignty + Jobs", "Building EU's First Music Tech Unicorn")

    eu_features = [
        ("🇪🇺 Digital Sovereignty", [
            "100% EU data infrastructure (Frankfurt/Amsterdam)",
            "Alternative to US-dominated market (DistroKid/TuneCore)",
            "All data in EU datacenters (Supabase EU)",
            "Zero transfers to non-EU jurisdictions"
        ]),
        ("💼 Job Creation Targets", [
            "Year 1-2: 15 direct hires (UK engineering, ops)",
            "Year 3-5: 50 hires (Berlin, Amsterdam offices)",
            "Year 5: 1,000+ sustainable creative jobs",
            "Indirect: 200+ jobs (engineers, curators)"
        ]),
        ("⚖️ Regulatory Leadership", [
            "First platform with full EU regulatory stack",
            "GDPR-native (privacy by design)",
            "DSA-ready (content moderation)",
            "EU AI Act pioneer (compliance model)"
        ]),
        ("🔬 Research Contribution", [
            "Open data API for academic research",
            "Training ground for EU AI talent (MCP dev)",
            "Climate-tech + music-tech convergence",
            "Built on European open-source (PostgreSQL)"
        ])
    ]

    idx = 0
    for row in range(2):
        for col in range(2):
            title, items = eu_features[idx]
            add_checklist_card(slide6, grid_x[col], grid_y[row], title, items)
            idx += 1

    # SLIDE 7: BUSINESS MODEL
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_gradient_background(slide7)
    add_slide_number(slide7, "(07)")
    add_title(slide7, "Business Model: €170M Revenue by Year 5", "Progressive Economics Aligned with Artist Success")

    # Commission flow
    commissions = [("20%", "Starting"), ("15%", "Growing"), ("10%", "Established"), ("2.5%", "Top Tier")]
    x_pos = 1.5
    for i, (percent, label) in enumerate(commissions):
        add_commission_tier(slide7, x_pos, 2.5, percent, label)
        x_pos += 3
        if i < 3:
            add_text_box(slide7, x_pos - 0.5, 2.7, 0.5, 0.5, "→", Pt(36), COLORS['light_gray'], False, PP_ALIGN.CENTER)

    # Revenue table
    revenue_items = [
        ("👥", "Artists", "500,000"),
        ("📊", "Annual Streams", "48 Billion"),
        ("💰", "Commission Revenue", "€48M"),
        ("🏢", "White-Label B2B", "€40M"),
        ("📱", "Subscriptions", "€60M"),
        ("⭐", "Premium Services", "€15M"),
        ("💎", "TOTAL REVENUE", "€170M")
    ]

    y_pos = 4.5
    for icon, label, value in revenue_items[:-1]:
        add_revenue_row(slide7, 2, y_pos, icon, label, value, False)
        y_pos += 0.5

    # Total row
    add_revenue_row(slide7, 2, y_pos, revenue_items[-1][0], revenue_items[-1][1], revenue_items[-1][2], True)

    # SLIDE 8: ADDITIONALITY
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_gradient_background(slide8)
    add_slide_number(slide8, "(08)")
    add_title(slide8, "Why EIC Grant? High-Risk R&D VCs Won't Fund", "Three Unsolved Problems Requiring Patient Capital")

    additionality = [
        ("🌍", "Public Good Innovation", "Carbon tracking methodology becomes industry standard • Open data API for research • Benefits entire music ecosystem, not just MSC & Co"),
        ("🔬", "Pre-Competitive R&D", "Sustainability measurement = unsolved problem • Blockchain adoption faces massive industry resistance • No proven commercial methodology exists"),
        ("📊", "Infrastructure-Heavy", "70%+ margins only at scale (Year 3+) • Sustainability compliance adds costs without immediate ROI • 24-month payback too long for early-stage VCs"),
        ("⚖️", "Regulatory Complexity", "GDPR/DSA/AI Act compliance requires legal resources • Multi-jurisdiction licensing (29 countries) • Carbon accounting standards still emerging")
    ]

    idx = 0
    for row in range(2):
        for col in range(2):
            icon, title, text = additionality[idx]
            add_feature_card(slide8, grid_x[col], grid_y[row], icon, title, text)
            idx += 1

    # SLIDE 9: USE OF FUNDS
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_gradient_background(slide9)
    add_slide_number(slide9, "(09)")
    add_title(slide9, "Use of Funds: €1.8M Over 24 Months", "Strategic Allocation Across Six Innovation Pillars")

    funds = [
        ("30%", "€540K", "🌍 Sustainability R&D", "• Carbon tracking engine (DIMPACT 2024)\n• Earth/Percent integration\n• Offset marketplace (5 providers)\n• Academic research partnerships"),
        ("20%", "€360K", "⛓️ Blockchain Infrastructure", "• Polygon smart contracts\n• Royalty verification system\n• NFT functionality\n• Copyright registration"),
        ("20%", "€360K", "🤖 AI & Automation", "• 181 MCP tools completion\n• Hit prediction + fraud detection\n• Apollo Intelligence enhancement\n• 15-language support"),
        ("15%", "€270K", "🌐 Distribution Partnerships", "• 150+ platform API integrations\n• Geographic expansion (Africa, LatAm)\n• White-label B2B development\n• Enterprise client onboarding"),
        ("10%", "€180K", "⚖️ Compliance & Legal", "• EU regulatory compliance\n• Patent applications (MCP + carbon)\n• B-Corp certification\n• Trademark registrations"),
        ("5%", "€90K", "🏢 Operations", "• Core team recruitment\n• Technical infrastructure\n• Contingency buffer\n• Office setup")
    ]

    grid_3x2_x = [0.8, 5.6, 10.4]
    grid_3x2_y = [2.5, 5.3]
    idx = 0
    for row in range(2):
        for col in range(3):
            percent, amount, title, items = funds[idx]
            add_fund_card(slide9, grid_3x2_x[col], grid_3x2_y[row], percent, amount, title, items)
            idx += 1

    # SLIDE 10: TEAM & ROADMAP
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_gradient_background(slide10)
    add_slide_number(slide10, "(10)")
    add_title(slide10, "Team & Roadmap: From 85% → Market Leader", "World-Class Team • Proven Track Record • Clear Execution Plan")

    # Team card
    team_text = """👑 Founder & CEO: Henry Taylor
14+ years gospel music industry experience • Organized Kingdom World Tour (Kirk Franklin, Maverick City Music) • BSc Biomedical Science + MRes Allied Medicine (LSE Distinction) • Full-stack developer (Next.js, React, PostgreSQL, AI integration) • Published world's first MCP server for music distribution • Festival of Life 2026 at The O2 Arena

Leadership Team Post-Funding (€280K Allocated):
💻 Feranmi Ogun - CTO & Technical Co-Founder
⚖️ Feyi Ogunseyinde - Head of Regulatory Compliance & Finance
🤝 Solomon Ade - Head of Commercial Partnerships & Distribution
🎨 Imisi Aina - VP of Product & Design
📢 Rebecca Taylor - Head of Marketing & Artist Growth
🚀 Opeyemi Kon Ajayi - Head of Engineering Operations"""

    add_team_card(slide10, 0.8, 2.3, 14.4, team_text)

    # Timeline
    timeline = [
        ("Q1 2026", "🚀 LAUNCH", "• Complete platform\n• 50-200 beta artists\n• Distribution APIs live"),
        ("Q2 2026", "5K ARTISTS", "• Public launch\n• Church network activated\n• £1M ARR achieved"),
        ("2027", "75K ARTISTS", "• €29M revenue\n• Series A funding secured\n• International expansion"),
        ("2030", "🏆 500K ARTISTS", "• €170M revenue\n• Carbon-neutral certified\n• B-Corp status\n• Market leadership")
    ]

    x_timeline = 0.8
    for period, milestone, details in timeline:
        add_timeline_card(slide10, x_timeline, 6, period, milestone, details)
        x_timeline += 3.7

    # Footer CTA
    add_impact_banner(slide10, 0.8, 8, 14.4,
                     "🚀 85% Complete Platform → 90-Day Launch → 18-24 Month Tech Lead",
                     "📧 info@htay.co.uk | 📱 +447530005553 | 🌐 staging.mscandco.com", COLORS['primary_blue'])

    return prs

# Helper functions
def add_dark_gradient_background(slide):
    """Add dark gradient background to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark']

def add_slide_number(slide, number):
    """Add slide number in top right"""
    num_box = slide.shapes.add_textbox(Inches(14.5), Inches(0.3), Inches(1.2), Inches(0.4))
    num_frame = num_box.text_frame
    num_frame.text = number
    p = num_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['light_gray']
    p.alignment = PP_ALIGN.RIGHT

def add_title(slide, title, subtitle):
    """Add title and subtitle to slide"""
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(14), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    subtitle_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(14), Inches(0.4))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['gold']

def add_innovation_card(slide, x, y, width, icon, text):
    """Add innovation card with icon and text"""
    card = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(2))
    card_frame = card.text_frame
    card_frame.text = f"{icon}\n{text}"

    # Icon paragraph
    p1 = card_frame.paragraphs[0]
    p1.font.size = Pt(42)
    p1.alignment = PP_ALIGN.CENTER

    # Text paragraph
    card_frame.add_paragraph()
    p2 = card_frame.paragraphs[1]
    p2.text = text
    p2.font.size = Pt(28)
    p2.font.color.rgb = COLORS['white']
    p2.alignment = PP_ALIGN.CENTER

def add_stat_card(slide, x, y, icon, value, label, sublabel):
    """Add stat card"""
    card = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(6.5), Inches(2.2))
    frame = card.text_frame
    frame.text = icon

    p1 = frame.paragraphs[0]
    p1.font.size = Pt(56)
    p1.alignment = PP_ALIGN.CENTER

    frame.add_paragraph()
    p2 = frame.paragraphs[1]
    p2.text = value
    p2.font.size = Pt(52)
    p2.font.bold = True
    p2.font.color.rgb = COLORS['primary_blue']
    p2.alignment = PP_ALIGN.CENTER

    frame.add_paragraph()
    p3 = frame.paragraphs[2]
    p3.text = label
    p3.font.size = Pt(22)
    p3.font.bold = True
    p3.font.color.rgb = COLORS['white']
    p3.alignment = PP_ALIGN.CENTER

    frame.add_paragraph()
    p4 = frame.paragraphs[3]
    p4.text = sublabel
    p4.font.size = Pt(18)
    p4.font.color.rgb = COLORS['light_gray']
    p4.alignment = PP_ALIGN.CENTER

def add_feature_card(slide, x, y, icon, title, text, width=6.5):
    """Add feature card"""
    card = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(2.2))
    frame = card.text_frame
    frame.text = icon

    p1 = frame.paragraphs[0]
    p1.font.size = Pt(60)
    p1.alignment = PP_ALIGN.CENTER

    frame.add_paragraph()
    p2 = frame.paragraphs[1]
    p2.text = title
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = COLORS['white']
    p2.alignment = PP_ALIGN.CENTER

    frame.add_paragraph()
    p3 = frame.paragraphs[2]
    p3.text = text
    p3.font.size = Pt(18)
    p3.font.color.rgb = COLORS['light_gray']
    p3.alignment = PP_ALIGN.LEFT

def add_comparison_text(slide, x, y, label, value, details, color):
    """Add comparison text section"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4), Inches(1))
    frame = box.text_frame
    frame.text = label

    p1 = frame.paragraphs[0]
    p1.font.size = Pt(18)
    p1.font.color.rgb = COLORS['light_gray']
    p1.alignment = PP_ALIGN.CENTER

    frame.add_paragraph()
    p2 = frame.paragraphs[1]
    p2.text = value
    p2.font.size = Pt(56)
    p2.font.bold = True
    p2.font.color.rgb = color
    p2.alignment = PP_ALIGN.CENTER

    frame.add_paragraph()
    p3 = frame.paragraphs[2]
    p3.text = details
    p3.font.size = Pt(16)
    p3.font.color.rgb = COLORS['light_gray']
    p3.alignment = PP_ALIGN.CENTER

def add_text_box(slide, x, y, width, height, text, font_size, color, bold, align):
    """Add simple text box"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.text = text
    p = frame.paragraphs[0]
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align

def add_impact_banner(slide, x, y, width, text, subtext, color):
    """Add impact banner"""
    banner = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.8))
    frame = banner.text_frame
    frame.text = text

    p1 = frame.paragraphs[0]
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = COLORS['white']
    p1.alignment = PP_ALIGN.CENTER

    if subtext:
        frame.add_paragraph()
        p2 = frame.paragraphs[1]
        p2.text = subtext
        p2.font.size = Pt(22)
        p2.font.color.rgb = COLORS['white']
        p2.alignment = PP_ALIGN.CENTER

def add_quote_banner(slide, x, y, width, quote):
    """Add quote banner"""
    banner = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.6))
    frame = banner.text_frame
    frame.text = quote
    p = frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.italic = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

def add_checklist_card(slide, x, y, title, items):
    """Add checklist card"""
    card = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(6.5), Inches(2.2))
    frame = card.text_frame
    frame.text = title

    p1 = frame.paragraphs[0]
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = COLORS['primary_blue']

    for item in items:
        frame.add_paragraph()
        p = frame.paragraphs[-1]
        p.text = f"✓ {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['white']
        p.level = 0

def add_commission_tier(slide, x, y, percent, label):
    """Add commission tier"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2), Inches(1))
    frame = box.text_frame
    frame.text = percent

    p1 = frame.paragraphs[0]
    p1.font.size = Pt(48)
    p1.font.bold = True
    p1.font.color.rgb = COLORS['primary_blue']
    p1.alignment = PP_ALIGN.CENTER

    frame.add_paragraph()
    p2 = frame.paragraphs[1]
    p2.text = label
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLORS['white']
    p2.alignment = PP_ALIGN.CENTER

def add_revenue_row(slide, x, y, icon, label, value, is_total):
    """Add revenue row"""
    row = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(12), Inches(0.4))
    frame = row.text_frame
    frame.text = f"{icon} {label}"

    p = frame.paragraphs[0]
    size = Pt(28) if is_total else Pt(20)
    p.font.size = size
    p.font.bold = is_total
    p.font.color.rgb = COLORS['white']

    # Value on right
    val_box = slide.shapes.add_textbox(Inches(x + 9), Inches(y), Inches(3), Inches(0.4))
    val_frame = val_box.text_frame
    val_frame.text = value
    p_val = val_frame.paragraphs[0]
    p_val.font.size = Pt(36) if is_total else Pt(20)
    p_val.font.bold = True
    p_val.font.color.rgb = COLORS['green'] if is_total else COLORS['white']
    p_val.alignment = PP_ALIGN.RIGHT

def add_fund_card(slide, x, y, percent, amount, title, items):
    """Add fund allocation card"""
    card = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4.5), Inches(2.5))
    frame = card.text_frame
    frame.text = f"{percent}  {amount}"

    p1 = frame.paragraphs[0]
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = COLORS['primary_blue']

    frame.add_paragraph()
    p2 = frame.paragraphs[1]
    p2.text = title
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = COLORS['white']

    frame.add_paragraph()
    p3 = frame.paragraphs[2]
    p3.text = items
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLORS['light_gray']

def add_team_card(slide, x, y, width, text):
    """Add team card"""
    card = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(3.3))
    frame = card.text_frame
    frame.text = text
    frame.word_wrap = True

    for p in frame.paragraphs:
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['dark']
        p.line_spacing = 1.3

def add_timeline_card(slide, x, y, period, milestone, details):
    """Add timeline card"""
    card = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(3.5), Inches(1.5))
    frame = card.text_frame
    frame.text = period

    p1 = frame.paragraphs[0]
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = COLORS['primary_blue']

    frame.add_paragraph()
    p2 = frame.paragraphs[1]
    p2.text = milestone
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = COLORS['white']

    frame.add_paragraph()
    p3 = frame.paragraphs[2]
    p3.text = details
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLORS['light_gray']

if __name__ == "__main__":
    print("Creating PowerPoint presentation...")
    prs = create_presentation()

    output_file = "MSC_EIC_Accelerator_10_Slides.pptx"
    prs.save(output_file)
    print(f"✅ PowerPoint created: {output_file}")

    print("\nTo convert to PDF, you can:")
    print("1. Open the .pptx file in PowerPoint/Keynote and export as PDF")
    print("2. Use an online converter")
    print("3. Use LibreOffice: soffice --headless --convert-to pdf *.pptx")
